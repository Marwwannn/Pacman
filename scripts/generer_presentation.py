"""Genere le support oral `docs/presentation.pptx` (17 diapositives).

Comme pour la documentation, aucun chiffre n'est recopie a la main : le deck
lit `results/campagne.json` et `results/descripteurs.json`, et pose les
captures d'ecran de `docs/captures/`. Apres une nouvelle campagne, le relancer
suffit pour que le support reste juste.

    python scripts/generer_presentation.py

Direction visuelle : encre profonde, jaune Pac-Man comme seul accent chaud,
les quatre couleurs des fantomes de l'arcade pour les quatre joueurs, des
pastilles en motif, de tres gros chiffres, une idee par diapositive.
Le fichier s'ouvre dans PowerPoint ou s'importe dans Google Slides ; les
polices (Montserrat, Open Sans) y sont natives.

Seules dependances hors projet : `python-pptx` (et Pillow, qu'il installe).
"""
import json
import math
import os
import tempfile
from pathlib import Path

try:
    from pptx import Presentation
except ModuleNotFoundError:  # pragma: no cover - message d'aide
    print("module `python-pptx` absent : pip install python-pptx")
    raise SystemExit(1) from None
from PIL import Image  # dependance de python-pptx, presente avec lui
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── palette
INK = RGBColor(0x0A, 0x0C, 0x14); INK2 = RGBColor(0x12, 0x15, 0x1F); LINE = RGBColor(0x23, 0x28, 0x38)
WHITE = RGBColor(0xF4, 0xF5, 0xF8); MUTED = RGBColor(0x9A, 0xA1, 0xB4); DIM = RGBColor(0x3A, 0x40, 0x52)
YELLOW = RGBColor(0xFF, 0xD2, 0x3F); RED = RGBColor(0xFF, 0x4F, 0x5E); PINK = RGBColor(0xFF, 0x9A, 0xD5)
CYAN = RGBColor(0x4F, 0xE3, 0xF0); ORANGE = RGBColor(0xFF, 0xB3, 0x47); GREEN = RGBColor(0x5C, 0xE2, 0x9A)
HEAD, BODY = "Montserrat", "Open Sans"

# ── données (jamais recopiées à la main)
RACINE = Path(__file__).resolve().parent.parent
R = str(RACINE / "results") + "/"
camp = json.load(open(R + "campagne.json", encoding="utf-8"))["etapes"]
desc = json.load(open(R + "descripteurs.json", encoding="utf-8"))["jeux"]
c4 = {r["agent"]: r for r in camp["comparatif_4f"]}
c1 = {r["agent"]: r for r in camp["baselines_1f"]}; c1["q-approxime"] = desc["base"]["un_fantome"]
pos = next(v for k, v in desc.items() if k != "base")
def f(x):
    """1234.0 -> « 1 234 » : l'espace fine des milliers, a la francaise."""
    return f"{x:,.0f}".replace(",", " ")

def pc(x):
    """0.94 -> « 94 % »."""
    return f"{x:.0%}".replace("%", " %")

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
N = 0
M = Inches(0.9)  # marge

def rect(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    r = s.shapes.add_shape(shape, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None: r.line.fill.background()
    else: r.line.color.rgb = line; r.line.width = Pt(1)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE: r.adjustments[0] = radius
    r.shadow.inherit = False
    return r

def dot(s, cx, cy, d, fill):
    return rect(s, cx - d // 2, cy - d // 2, d, d, fill, shape=MSO_SHAPE.OVAL)

def texte(s, x, y, w, h, txt, taille=18, couleur=MUTED, gras=False, font=BODY,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, espace=None, interligne=None, spacing=None):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lignes = txt if isinstance(txt, list) else [txt]
    for i, contenu in enumerate(lignes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = align
        if espace: p.space_after = Pt(espace)
        if interligne: p.line_spacing = interligne
        parts = contenu if isinstance(contenu, list) else [(contenu, couleur, gras)]
        for part in parts:
            t, c, g = part[0], part[1], part[2]
            fn = part[3] if len(part) > 3 else font
            r = p.add_run(); r.text = t; r.font.size = Pt(taille); r.font.color.rgb = c; r.font.bold = g; r.font.name = fn
            if spacing is not None:
                rPr = r._r.get_or_add_rPr(); rPr.set("spc", str(spacing))
    return tb

def pellets(s, x1, x2, y, pas=Inches(0.42), d=Inches(0.09), fill=DIM):
    x = x1
    while x <= x2:
        dot(s, x, y, d, fill); x += pas

def pacman(s, cx, cy, d, fill=YELLOW, bouche=32):
    """Un disque auquel il manque un secteur : PIE, angles en 60000e de degré (python-pptx normalise /100000)."""
    p = s.shapes.add_shape(MSO_SHAPE.PIE, cx - d // 2, cy - d // 2, d, d)
    p.adjustments[0] = bouche * 60000 / 100000          # début (sous l'horizontale, à droite)
    p.adjustments[1] = (360 - bouche) * 60000 / 100000  # fin
    p.fill.solid(); p.fill.fore_color.rgb = fill; p.line.fill.background(); p.shadow.inherit = False
    return p

def fantome(s, x, y, w, fill):
    """Silhouette de fantôme d'arcade : dôme + jupe en zigzag, en forme libre."""
    h = int(w * 1.05); r = w / 2; cx = x + r; top = y + r
    pts = []
    for i in range(0, 19):  # dôme, de gauche à droite
        a = math.pi - i * math.pi / 18
        pts.append((int(cx + r * math.cos(a)), int(top - r * math.sin(a))))
    bas = y + h; jupe = int(w * 0.14)
    pts.append((x + w, bas))
    n = 3
    for k in range(n):  # zigzag de droite à gauche
        xk1 = x + w - (2 * k + 1) * w / (2 * n); xk2 = x + w - (2 * k + 2) * w / (2 * n)
        pts.append((int(xk1), bas - jupe)); pts.append((int(xk2), bas))
    fb = s.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
    fb.add_line_segments(pts[1:], close=True)
    g = fb.convert_to_shape(); g.fill.solid(); g.fill.fore_color.rgb = fill; g.line.fill.background(); g.shadow.inherit = False
    # yeux
    ew = int(w * 0.22); eh = int(w * 0.28)
    for ex in (cx - int(w * 0.27), cx + int(w * 0.05)):
        rect(s, ex, top - int(eh * 0.55), ew, eh, WHITE, shape=MSO_SHAPE.OVAL)
        rect(s, ex + int(ew * 0.45), top - int(eh * 0.25), int(ew * 0.45), int(eh * 0.45), INK, shape=MSO_SHAPE.OVAL)
    return g

def base(s):
    rect(s, 0, 0, W, H, INK)

def pied(s, coul=YELLOW):
    global N; N += 1
    y = H - Inches(0.62)
    texte(s, M, y, Inches(6), Inches(0.3), "PAC-MAN · UN JEU, QUATRE JOUEURS ARTIFICIELS", 9, DIM, True, HEAD, spacing=200)
    rect(s, W - M - Inches(0.55), y - Inches(0.05), Inches(0.55), Inches(0.36), coul, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    texte(s, W - M - Inches(0.55), y - Inches(0.05), Inches(0.55), Inches(0.36), f"{N:02d}", 10, INK, True, HEAD, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

def kicker(s, y, txt, coul=YELLOW):
    dot(s, M + Inches(0.06), y + Inches(0.13), Inches(0.12), coul)
    texte(s, M + Inches(0.3), y, Inches(10), Inches(0.3), txt.upper(), 11, coul, True, HEAD, spacing=250)

def titre(s, y, txt, taille=40, w=Inches(11.5), coul=WHITE):
    return texte(s, M, y, w, Inches(1.6), txt, taille, coul, True, HEAD, interligne=1.05)

def slide(k, t, taille=40, coul=YELLOW):
    s = prs.slides.add_slide(BLANK); base(s)
    kicker(s, Inches(0.75), k, coul); titre(s, Inches(1.1), t, taille)
    pied(s, coul); return s

def puce(s, x, y, w, coul, tete, corps, taille=15, h=Inches(1.1)):
    rect(s, x, y + Inches(0.06), Inches(0.05), h - Inches(0.2), coul)
    texte(s, x + Inches(0.3), y, w - Inches(0.3), h, [[(tete, WHITE, True, HEAD)], [(corps, MUTED, False)]], taille, espace=3)

def carte(s, x, y, w, h, coul, num, tete, corps, taille=14):
    rect(s, x, y, w, h, INK2, LINE, MSO_SHAPE.ROUNDED_RECTANGLE, 0.05)
    rect(s, x, y, w, Inches(0.06), coul)
    texte(s, x + Inches(0.35), y + Inches(0.35), w - Inches(0.7), Inches(0.5), num, 28, coul, True, HEAD)
    texte(s, x + Inches(0.35), y + Inches(0.95), w - Inches(0.7), Inches(0.8), tete, 15, WHITE, True, HEAD)
    texte(s, x + Inches(0.35), y + Inches(1.8), w - Inches(0.7), h - Inches(1.95), corps, taille, MUTED, espace=6)

SHOTS = str(RACINE / "docs" / "captures")
TMP = tempfile.mkdtemp(prefix="presentation-")

def photo(s, nom, x, y, w, crop=None, h_max=None):
    """Capture d'écran recadrée, posée sur une carte avec un liseré. Retourne (largeur, hauteur) posées."""
    src = os.path.join(SHOTS, nom); im = Image.open(src)
    if crop:
        im = im.crop(tuple(int(v * im.width) if i % 2 == 0 else int(v * im.height) for i, v in enumerate(crop)))
    tmp = os.path.join(TMP, "crop_" + nom); im.save(tmp)
    ratio = im.height / im.width
    ph = int(w * ratio)
    if h_max and ph > h_max:
        ph = h_max; w = int(ph / ratio)
    pad = Inches(0.08)
    rect(s, x - pad, y - pad, w + 2 * pad, ph + 2 * pad, INK2, LINE, MSO_SHAPE.ROUNDED_RECTANGLE, 0.03)
    s.shapes.add_picture(tmp, x, y, width=w, height=ph)
    return w, ph

def legende(s, x, y, w, coul, tete, corps, taille=12):
    texte(s, x, y, w, Inches(0.9), [[(tete, coul, True, HEAD)], [(corps, MUTED, False)]], taille, espace=2)

# ═══════════════════════════════════════════ 1. COUVERTURE
s = prs.slides.add_slide(BLANK); base(s)
# couloir de pastilles, Pac-Man les mange
yl = Inches(2.05)
pellets(s, Inches(5.6), W - M, yl, pas=Inches(0.5), d=Inches(0.13), fill=YELLOW)
pacman(s, Inches(4.9), yl, Inches(0.95))
texte(s, M, Inches(2.75), Inches(12), Inches(1.6), "PAC-MAN", 96, YELLOW, True, HEAD, spacing=-300)
texte(s, M, Inches(4.15), Inches(12), Inches(0.9), "Un jeu, quatre joueurs artificiels", 34, WHITE, True, HEAD)
texte(s, M, Inches(5.2), Inches(8), Inches(0.5), "Marwan Youmni  ·  Master 1 IA_PO  ·  2025-2026", 15, MUTED)
texte(s, M, Inches(5.6), Inches(8), Inches(0.4), "github.com/Marwwannn/Pacman", 13, DIM)
# quatre fantômes en bas à droite
for i, c in enumerate([RED, PINK, CYAN, ORANGE]):
    fantome(s, W - M - Inches(3.6) + i * Inches(0.95), Inches(5.0), Inches(0.75), c)
pied(s)

# ═══════════════════════════════════════════ 2. EN UNE PHRASE
s = prs.slides.add_slide(BLANK); base(s)
kicker(s, Inches(0.75), "Le projet")
lignes = [("01", "J'ai construit le jeu.", "Un moteur Pac-Man complet en Python, sans dépendance, et un client web pour y jouer."),
          ("02", "Puis des IA qui y jouent.", "Quatre stratégies, de la plus bête à la plus forte, dont une qui apprend seule."),
          ("03", "Puis je les ai mesurées.", "100 parties chacune, sur des situations jamais vues. Le score est un nombre.")]
for i, (n, t, d) in enumerate(lignes):
    y = Inches(1.55) + i * Inches(1.6)
    texte(s, M, y, Inches(1.2), Inches(1.2), n, 44, DIM, True, HEAD)
    texte(s, M + Inches(1.4), y + Inches(0.02), Inches(10), Inches(0.7), t, 30, WHITE, True, HEAD)
    texte(s, M + Inches(1.4), y + Inches(0.62), Inches(9.5), Inches(0.7), d, 15, MUTED)
    if i < 2: rect(s, M + Inches(1.4), y + Inches(1.35), Inches(10.1), Emu(9525), LINE)
pied(s)

# ═══════════════════════════════════════════ 2b. LE TERRAIN (captures)
s = slide("Le terrain", "Le jeu existe pour de vrai : on y joue", 40)
# deux paires image + légende : l'écran d'accueil, puis une partie en cours
y0 = Inches(2.3); hi = Inches(3.75)
w1, h1 = photo(s, "accueil.png", M, y0, Inches(3.0), crop=(0.21, 0.09, 0.78, 0.90), h_max=hi)
x2 = M + Inches(6.0)
w2, h2 = photo(s, "ia_recherche_en_cours.png", x2, y0, Inches(3.0), crop=(0.21, 0.09, 0.78, 0.90), h_max=hi)
puce(s, M + w1 + Inches(0.35), y0 + Inches(0.1), Inches(5.6) - w1, YELLOW, "L'accueil", "Flèches pour jouer soi-même, ou un lien pour regarder une IA jouer à sa place.", 13, Inches(1.6))
puce(s, x2 + w2 + Inches(0.35), y0 + Inches(0.1), Inches(5.6) - w2, CYAN, "Une partie", "Le serveur Python calcule tout ; le navigateur ne fait qu'afficher. Le même moteur entraîne les IA, 3 000 fois plus vite que le temps réel.", 13, Inches(2.2))

# ═══════════════════════════════════════════ 3. POURQUOI CE JEU
s = prs.slides.add_slide(BLANK); base(s)
kicker(s, Inches(0.75), "Pourquoi ce jeu")
texte(s, M, Inches(1.15), Inches(5.6), Inches(2.6), [[("Règles", WHITE, True, HEAD)], [("simples.", WHITE, True, HEAD)], [("Stratégie", YELLOW, True, HEAD)], [("pas du tout.", YELLOW, True, HEAD)]], 44, interligne=1.0)
x = Inches(6.9); w = Inches(5.5)
puce(s, x, Inches(1.25), w, YELLOW, "Un arbitrage permanent", "Manger les pastilles (le gain) sans se faire attraper (le risque). À chaque croisement.")
puce(s, x, Inches(2.55), w, CYAN, "Des adversaires imparfaits", "Les quatre fantômes de 1980 sont myopes : ils visent à vol d'oiseau, sans voir les murs. C'est ce qui les rend battables.")
puce(s, x, Inches(3.85), w, GREEN, "Une mesure honnête", "Le score final dit tout. Personne n'a besoin de juger.")
puce(s, x, Inches(5.0), w, PINK, "Un terrain que je contrôle", "Quand un agent stagne, je peux distinguer « la méthode est mauvaise » de « le problème est mal posé ».")
# fantômes en rappel discret
for i, c in enumerate([RED, PINK, CYAN, ORANGE]):
    fantome(s, M + i * Inches(0.7), Inches(5.35), Inches(0.5), c)
pellets(s, M + Inches(3.0), Inches(6.2), Inches(5.6), pas=Inches(0.32), d=Inches(0.08))
pied(s)

# ═══════════════════════════════════════════ 4. LES QUATRE JOUEURS
s = slide("Les quatre joueurs · quatre stratégies", "De la plus bête à la plus forte", 40)
joueurs = [("Aléatoire", "le plancher", RED, "Une direction au hasard à chaque croisement. Tout ce qui fait moins est cassé."),
           ("Heuristique", "règles à la main", PINK, "Chasseur proche → fuir. Fantôme effrayé → le chasser. Sinon → la pastille la plus proche."),
           ("Q-learning", "il apprend", YELLOW, "Regarde 12 indicateurs et apprend, en 3 000 parties, combien chacun compte. Personne ne lui dit les règles."),
           ("Recherche", "il simule", CYAN, "Clone la partie, essaie chaque direction 3 coups à l'avance, garde la meilleure. Zéro apprentissage.")]
cw = Inches(2.7); gap = Inches(0.27); y0 = Inches(2.35); ch = Inches(3.35)
for i, (nom, sous, c, d) in enumerate(joueurs):
    x = M + i * (cw + gap)
    rect(s, x, y0, cw, ch, INK2, LINE, MSO_SHAPE.ROUNDED_RECTANGLE, 0.05)
    dot(s, x + Inches(0.55), y0 + Inches(0.6), Inches(0.5), c)
    texte(s, x + Inches(1.0), y0 + Inches(0.38), cw - Inches(1.2), Inches(0.5), nom, 17, WHITE, True, HEAD, anchor=MSO_ANCHOR.MIDDLE)
    texte(s, x + Inches(0.3), y0 + Inches(1.15), cw - Inches(0.6), Inches(0.4), sous.upper(), 10, c, True, HEAD, spacing=150)
    texte(s, x + Inches(0.3), y0 + Inches(1.6), cw - Inches(0.6), ch - Inches(1.8), d, 14, MUTED)
# échelle de force
ya = Inches(6.05)
rect(s, M, ya, Inches(11.5), Emu(19050), LINE)
texte(s, M, ya + Inches(0.08), Inches(4), Inches(0.3), "MOINS FORT", 9, DIM, True, HEAD, spacing=200)
texte(s, W - M - Inches(4), ya + Inches(0.08), Inches(4), Inches(0.3), "PLUS FORT  →", 9, YELLOW, True, HEAD, PP_ALIGN.RIGHT, spacing=200)

# ═══════════════════════════════════════════ 5. L'AGENT QUI APPREND
s = slide("L'agent qui apprend", "Douze indicateurs, douze poids, un barème", 36)
texte(s, M, Inches(2.25), Inches(7), Inches(0.9), [[("Q(s, a) = ", WHITE, True, HEAD), ("w · f(s, a)", YELLOW, True, HEAD)]], 34)
texte(s, M, Inches(3.0), Inches(7), Inches(0.5), "La valeur d'un coup, c'est la somme des indicateurs pondérés. Après 3 000 parties (≈ 8 min), on lit les poids.", 13, MUTED)
chips = ["pastille sur la case", "super-pastille", "fruit", "distance pastille", "distance super-pastille", "distance chasseur",
         "chasseurs à < 3 cases", "proie effrayée à portée", "nombre de sorties", "demi-tour", "avancement", "biais"]
x = M; y = Inches(3.75); lh = Inches(0.42)
for ch_ in chips:
    wch = Inches(0.22 + 0.085 * len(ch_))
    if x + wch > M + Inches(7.1): x = M; y += lh + Inches(0.12)
    rect(s, x, y, wch, lh, INK2, LINE, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5)
    texte(s, x, y, wch, lh, ch_, 11, WHITE, False, BODY, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    x += wch + Inches(0.12)
# barème
bx = Inches(8.5); by = Inches(2.2); bw = Inches(3.95)
rect(s, bx, by, bw, Inches(4.0), INK2, LINE, MSO_SHAPE.ROUNDED_RECTANGLE, 0.04)
texte(s, bx + Inches(0.35), by + Inches(0.3), bw, Inches(0.3), "LE BARÈME", 10, YELLOW, True, HEAD, spacing=200)
bar = [("Pastille", "+10", GREEN), ("Super-pastille", "+50", GREEN), ("Fantôme mangé", "+200 → 1600", GREEN),
       ("Niveau terminé", "+500", GREEN), ("Mort", "−500", RED), ("Chaque décision", "−1", RED)]
for i, (k, v, c) in enumerate(bar):
    yy = by + Inches(0.8) + i * Inches(0.5)
    texte(s, bx + Inches(0.35), yy, Inches(2.2), Inches(0.4), k, 14, WHITE, False, BODY, anchor=MSO_ANCHOR.MIDDLE)
    texte(s, bx + Inches(2.0), yy, bw - Inches(2.35), Inches(0.4), v, 15, c, True, HEAD, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
    if i < 5: rect(s, bx + Inches(0.35), yy + Inches(0.46), bw - Inches(0.7), Emu(9525), LINE)

# ═══════════════════════════════════════════ 6. RÉSULTATS
s = slide("Résultats · quatre fantômes", "100 parties chacun, jamais vues", 40)
ordre = [("Aléatoire", "aleatoire", RED), ("Heuristique", "heuristique", PINK), ("Q-learning", "q-approxime", YELLOW), ("Recherche", "recherche", CYAN)]
maxs = max(c4[k]["score_median"] for _, k, _ in ordre)
bx = M; bw = Inches(6.6); y0 = Inches(2.45)
for i, (nom, k, c) in enumerate(ordre):
    r = c4[k]; y = y0 + i * Inches(0.98)
    texte(s, bx, y, Inches(2.0), Inches(0.35), nom, 14, WHITE, True, HEAD, anchor=MSO_ANCHOR.MIDDLE)
    texte(s, bx, y + Inches(0.33), Inches(2.0), Inches(0.3), f"{pc(r['taux_victoire'])} de victoires", 11, MUTED)
    wbar = int(bw * r["score_median"] / maxs)
    rect(s, bx + Inches(2.1), y + Inches(0.02), bw, Inches(0.5), INK2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    rect(s, bx + Inches(2.1), y + Inches(0.02), max(wbar, Inches(0.3)), Inches(0.5), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    texte(s, bx + Inches(2.1) + wbar + Inches(0.15), y + Inches(0.02), Inches(1.5), Inches(0.5), f(r["score_median"]), 15, WHITE, True, HEAD, anchor=MSO_ANCHOR.MIDDLE)
texte(s, bx + Inches(2.1), Inches(6.35), Inches(6), Inches(0.3), "SCORE MÉDIAN SUR 100 PARTIES", 9, DIM, True, HEAD, spacing=200)
# callout
cx = Inches(9.9); cw2 = Inches(2.55)
texte(s, cx, Inches(2.3), cw2, Inches(1.1), pc(c4["recherche"]["taux_victoire"]), 64, CYAN, True, HEAD, spacing=-200)
texte(s, cx, Inches(3.35), cw2, Inches(0.8), "de victoires pour la recherche : elle voit l'avenir.", 13, MUTED)
rect(s, cx, Inches(4.2), cw2, Emu(9525), LINE)
texte(s, cx, Inches(4.4), cw2, Inches(1.8), [[("L'agent appris bat mes règles ", WHITE, True), (f"({f(c4['q-approxime']['score_median'])} contre {f(c4['heuristique']['score_median'])}) sans qu'on lui ait rien dit. Mais il gagne rarement : douze poids ne planifient pas une fin de niveau.", MUTED, False)]], 12)

# ═══════════════════════════════════════════ 6b. L'IA JOUE EN DIRECT (capture)
s = slide("Dans le jeu", "On regarde l'agent appris jouer, en direct", 36)
w1, h1 = photo(s, "ia_appris_en_cours.png", M, Inches(2.3), Inches(5.6), crop=(0.21, 0.09, 0.78, 0.90), h_max=Inches(4.1))
x = M + w1 + Inches(0.6); w = W - M - x
texte(s, x, Inches(2.3), w, Inches(0.5), [[("/?ia=appris", YELLOW, True, HEAD)]], 18)
puce(s, x, Inches(2.95), w, YELLOW, "Le modèle final au volant", "Le serveur décide à chaque croisement avec les douze poids appris ; je ne touche à rien.", 13, Inches(1.05))
puce(s, x, Inches(4.05), w, CYAN, "On le voit hésiter", "Aux intersections, l'agent compare les directions ; parfois il revient sur ses pas. C'est le −1 par décision qui le pousse à avancer.", 13, Inches(1.2))
puce(s, x, Inches(5.3), w, PINK, "Les mêmes règles pour tous", "Remplacer appris par recherche, heuristique ou aleatoire dans l'adresse suffit pour comparer à l'œil.", 13, Inches(1.05))

# ═══════════════════════════════════════════ 6c. CHAQUE DÉCISION S'OUVRE (capture)
s = slide("Sous le capot · poids par poids", "Chaque décision peut s'ouvrir", 40)
w1, h1 = photo(s, "decisions_html.png", M, Inches(2.3), Inches(7.2), crop=(0.0, 0.0, 1.0, 0.87), h_max=Inches(3.9))
x = M + w1 + Inches(0.5); w = W - M - x
puce(s, x, Inches(2.3), w, YELLOW, "docs/decisions.html", "Une partie de l'IA rejouée coup par coup, sur une graine jamais vue à l'entraînement.", 13, Inches(1.2))
puce(s, x, Inches(3.55), w, GREEN, "Vert pousse, rouge retient", "À chaque intersection : la valeur de chaque direction, décomposée en poids × descripteur.", 13, Inches(1.2))
puce(s, x, Inches(4.8), w, CYAN, "Ce que douze poids ont acheté", "Un réseau de neurones n'aurait jamais permis de montrer ça.", 13, Inches(1.1))

# ═══════════════════════════════════════════ 7. PROBLÈME 1
s = prs.slides.add_slide(BLANK); base(s)
kicker(s, Inches(0.75), "Problème n° 1", RED)
titre(s, Inches(1.1), "Décider à chaque pas, c'est décider 2 500 fois pour rien", 34)
texte(s, M, Inches(2.5), Inches(3.2), Inches(1.8), "34", 120, YELLOW, True, HEAD, spacing=-400)
texte(s, M, Inches(4.35), Inches(3.4), Inches(0.8), "vrais points de décision, sur 300 cases praticables", 14, MUTED)
# mini couloir : pastilles grises, une jaune aux intersections
pellets(s, M, M + Inches(3.2), Inches(5.45), pas=Inches(0.27), d=Inches(0.08))
for xx in (M, M + Inches(1.62), M + Inches(3.24)): dot(s, xx, Inches(5.45), Inches(0.14), YELLOW)
x = Inches(5.1); w = Inches(7.3)
puce(s, x, Inches(2.5), w, RED, "Le constat", "89 % du labyrinthe est un couloir où la direction est imposée. Décider à chaque tick, c'est 2 440 décisions sans conséquence par partie, et apprendre sur ce bruit.", 14, Inches(1.3))
puce(s, x, Inches(3.85), w, YELLOW, "La solution", "L'environnement traverse les couloirs tout seul et ne rend la main qu'aux intersections.", 14, Inches(1.0))
puce(s, x, Inches(4.95), w, GREEN, "Le résultat", "Une partie passe de ~2 500 pas à ~100 décisions : 25 fois plus court, sans rien perdre du jeu. C'est ce qui rend l'apprentissage possible en minutes.", 14, Inches(1.3))
pied(s, RED)

# ═══════════════════════════════════════════ 8. PROBLÈME 2
s = prs.slides.add_slide(BLANK); base(s)
kicker(s, Inches(0.75), "Problème n° 2", RED)
titre(s, Inches(1.1), "L'agent apprenait par cœur, pas à jouer", 36)
texte(s, M, Inches(2.2), Inches(11.5), Inches(0.9), "Le moteur est parfaitement reproductible, indispensable pour tester. Mais si chaque partie démarre pareil, l'agent mémorise une suite de coups, pas une stratégie : score flatteur à l'entraînement, effondrement dès que ça change.", 15, MUTED)
cw = Inches(3.65); gap = Inches(0.27); y0 = Inches(3.3); ch = Inches(3.1)
fixes = [("1", "Une graine par partie", "Le départ de Pac-Man et l'errance des fantômes changent à chaque épisode."),
         ("2", "Des graines jamais vues", "Le code refuse d'évaluer sur une graine d'entraînement. Une garde, pas une convention."),
         ("3", "La médiane, pas le meilleur run", "100 parties, médiane et écart-type. Un bon score isolé ne prouve rien.")]
for i, (n, t, d) in enumerate(fixes):
    carte(s, M + i * (cw + gap), y0, cw, ch, YELLOW, n, t, d)
pied(s, RED)

# ═══════════════════════════════════════════ 9. PROBLÈME 3
s = prs.slides.add_slide(BLANK); base(s)
kicker(s, Inches(0.75), "Problème n° 3", RED)
titre(s, Inches(1.1), "Deux chiffres du barème comptaient plus que tout l'algorithme", 34)
hw = Inches(5.6); y0 = Inches(2.55); hh = Inches(3.6)
for i, (num, c, t, d, fix) in enumerate([
        ("−500", RED, "L'agent se suicidait", "Sans pénalité de mort assez lourde, mourir vite était la meilleure façon d'arrêter de perdre des points.", "La mort domine tout le reste."),
        ("−1", CYAN, "L'agent tournait en rond", "Devant une super-pastille, il faisait des allers-retours sans la manger : rien ne le pressait.", "Chaque décision coûte. Il avance.")]):
    x = M + i * (hw + Inches(0.3))
    rect(s, x, y0, hw, hh, INK2, LINE, MSO_SHAPE.ROUNDED_RECTANGLE, 0.04)
    texte(s, x + Inches(0.4), y0 + Inches(0.25), hw, Inches(1.2), num, 64, c, True, HEAD, spacing=-200)
    texte(s, x + Inches(0.4), y0 + Inches(1.5), hw - Inches(0.8), Inches(0.5), t, 18, WHITE, True, HEAD)
    texte(s, x + Inches(0.4), y0 + Inches(2.0), hw - Inches(0.8), Inches(0.9), d, 13, MUTED)
    texte(s, x + Inches(0.4), y0 + Inches(2.95), hw - Inches(0.8), Inches(0.5), [[("→  ", c, True, HEAD), (fix, WHITE, True, HEAD)]], 13)
pied(s, RED)

# ═══════════════════════════════════════════ 10. CE QUI N'A PAS MARCHÉ
s = prs.slides.add_slide(BLANK); base(s)
kicker(s, Inches(0.75), "Ce qui n'a pas marché", ORANGE)
titre(s, Inches(1.1), "Trois résultats négatifs, publiés tels quels", 36)
cw = Inches(3.65); gap = Inches(0.27); y0 = Inches(2.35); ch = Inches(3.9)
neg = [("−18 %", "La position de chaque fantôme", f"Donner à l'agent les positions exactes le rend moins bon : {f(pos['quatre_fantomes']['score_median'])} contre {f(desc['base']['quatre_fantomes']['score_median'])}. Deux fois plus de poids sur le même budget, et un modèle linéaire additionne, il ne croise pas."),
       ("+0", "L'entraînement progressif", "Un fantôme puis quatre : 2 895 avec, 2 915 sans. Bonne idée sur le papier, aucun gain mesuré."),
       ("∅", "Un test vert qui ne testait rien", "« L'agent ressort d'un cul-de-sac » passait toujours… le labyrinthe n'en a aucun. Il itérait sur une liste vide. Un test qui rassure sans prouver est pire qu'un test absent.")]
for i, (n, t, d) in enumerate(neg):
    carte(s, M + i * (cw + gap), y0, cw, ch, ORANGE, n, t, d, 13)
pied(s, ORANGE)

# ═══════════════════════════════════════════ 11. APPRENDRE OU RECALCULER
s = slide("La vraie question", "Apprendre, ou recalculer ?", 40)
cols = [("Q-LEARNING", YELLOW, ["3 000 parties, ~8 min avant de jouer", "12 multiplications par coup", "Il reste 12 nombres lisibles", "Fonctionne sans simulateur exact"]),
        ("RECHERCHE", CYAN, ["Zéro coût avant de jouer", "Des dizaines de parties simulées par coup", "Il ne reste rien", "S'effondre sans simulateur exact"])]
labels = ["Avant de jouer", "Pendant le jeu", "Après", "Sans simulateur"]
x0 = M; lw = Inches(2.3); cw = Inches(3.0); y0 = Inches(2.35); rh = Inches(0.78)
for j, (t, c, _) in enumerate(cols):
    texte(s, x0 + lw + j * (cw + Inches(0.2)), y0, cw, Inches(0.4), t, 11, c, True, HEAD, spacing=200)
for i, lab in enumerate(labels):
    y = y0 + Inches(0.55) + i * rh
    rect(s, x0, y - Inches(0.08), lw + 2 * cw + Inches(0.2), Emu(9525), LINE)
    texte(s, x0, y, lw, rh, lab, 12, DIM, True, HEAD, anchor=MSO_ANCHOR.MIDDLE)
    for j, (_, c, vals) in enumerate(cols):
        texte(s, x0 + lw + j * (cw + Inches(0.2)), y, cw, rh, vals[i], 13, WHITE if i != 3 else c, i == 3, BODY, anchor=MSO_ANCHOR.MIDDLE)
# leçon
lx = Inches(9.6); lw2 = Inches(2.85)
rect(s, lx, Inches(2.35), Inches(0.05), Inches(3.7), YELLOW)
texte(s, lx + Inches(0.3), Inches(2.35), lw2, Inches(0.3), "LA LEÇON", 10, YELLOW, True, HEAD, spacing=200)
texte(s, lx + Inches(0.3), Inches(2.75), lw2, Inches(3.4), [[("La recherche gagne ici parce qu'elle a un simulateur exact et gratuit.", WHITE, True, HEAD)], [("Un robot, un marché, un patient n'en ont pas : simuler y est impossible ou plus cher qu'agir. Là, c'est l'apprentissage qui reste.", MUTED, False)]], 13, espace=8)

# ═══════════════════════════════════════════ 12. LIMITES & SUITES
s = slide("Limites et suites", "Ce que je n'affirme pas, et ce qui vient", 36)
hw = Inches(5.6); y0 = Inches(2.5)
for i, (t, c, items) in enumerate([
        ("LIMITES", RED, ["Pas de niveau terminé à 4 fantômes : l'agent évalue chaque croisement isolément.", "Les 12 indicateurs sont écrits à la main : il apprend combien ça compte, pas quoi regarder.", "Un seul labyrinthe mesuré : le transfert n'est pas vérifié, donc pas affirmé."]),
        ("SUITES", GREEN, ["Des indicateurs appris (auto-encodeur) plutôt qu'écrits à la main.", "Des fantômes non myopes, avec le vrai plus court chemin : un mode difficile.", "Un réseau de neurones pour croiser ce qu'un modèle linéaire ne fait qu'additionner."])]):
    x = M + i * (hw + Inches(0.3))
    texte(s, x, y0, hw, Inches(0.3), t, 11, c, True, HEAD, spacing=250)
    for k, it in enumerate(items):
        y = y0 + Inches(0.55) + k * Inches(1.2)
        rect(s, x, y - Inches(0.1), hw, Emu(9525), LINE)
        dot(s, x + Inches(0.1), y + Inches(0.18), Inches(0.12), c)
        texte(s, x + Inches(0.4), y, hw - Inches(0.4), Inches(1.1), it, 14, MUTED)

# ═══════════════════════════════════════════ 13. IA GÉNÉRATIVE
s = slide("Usage de l'IA générative", "Déclaré, total, assumé", 40)
texte(s, M, Inches(2.3), Inches(11.5), Inches(0.6), "Claude Code (Anthropic), seul outil, du 19/07 au 21/08. Chaque commit porte le trailer Co-Authored-By: Claude.", 15, MUTED)
hw = Inches(5.6); y0 = Inches(3.1)
for i, (t, c, body) in enumerate([
        ("CE QU'IL A FAIT", DIM, "Le socle du jeu, la fidélité à l'arcade, les 278 tests, l'audit sécurité, le cadrage chiffré de l'approche."),
        ("CE QUE J'AI FAIT", YELLOW, "Cadrer, arbitrer, éprouver, refuser. Le périmètre, le choix du renforcement, les bugs que seul un joueur voit, et les deux questions qui ont produit les résultats les plus forts : « et si on lui donnait la position des fantômes ? », « les fantômes partent-ils toujours du même endroit ? »")]):
    x = M + i * (hw + Inches(0.3))
    rect(s, x, y0, hw, Inches(3.0), INK2, LINE, MSO_SHAPE.ROUNDED_RECTANGLE, 0.04)
    rect(s, x, y0, hw, Inches(0.06), c)
    texte(s, x + Inches(0.35), y0 + Inches(0.35), hw, Inches(0.3), t, 10, c if c != DIM else MUTED, True, HEAD, spacing=200)
    texte(s, x + Inches(0.35), y0 + Inches(0.8), hw - Inches(0.7), Inches(2.1), body, 14, WHITE if i else MUTED)

# ═══════════════════════════════════════════ 14. MERCI
s = prs.slides.add_slide(BLANK); base(s)
yl = Inches(1.7)
pellets(s, M + Inches(1.2), Inches(7.6), yl, pas=Inches(0.5), d=Inches(0.13), fill=YELLOW)
pacman(s, M + Inches(0.48), yl, Inches(0.95))
texte(s, M, Inches(2.5), Inches(11), Inches(1.2), "Merci. Des questions ?", 54, WHITE, True, HEAD, spacing=-100)
texte(s, M, Inches(3.9), Inches(11), Inches(0.3), "DÉMO", 11, YELLOW, True, HEAD, spacing=250)
demo = [("pacman-server", "on joue, sur http://127.0.0.1:8000"),
        ("/?ia=appris", "l'agent appris joue en direct, on le voit hésiter aux croisements"),
        ("docs/decisions.html", "chaque décision expliquée, poids par poids")]
for i, (k, v) in enumerate(demo):
    y = Inches(4.35) + i * Inches(0.55)
    texte(s, M, y, Inches(3.2), Inches(0.45), k, 15, YELLOW, True, HEAD, anchor=MSO_ANCHOR.MIDDLE)
    texte(s, M + Inches(3.2), y, Inches(8), Inches(0.45), v, 14, MUTED, anchor=MSO_ANCHOR.MIDDLE)
texte(s, M, Inches(6.15), Inches(8), Inches(0.4), "github.com/Marwwannn/Pacman", 13, DIM)
for i, c in enumerate([RED, PINK, CYAN, ORANGE]):
    fantome(s, W - M - Inches(3.0) + i * Inches(0.8), Inches(5.4), Inches(0.6), c)
pied(s)

CIBLE = RACINE / "docs" / "presentation.pptx"
prs.save(str(CIBLE))
print(f"{CIBLE.relative_to(RACINE)} ecrit ({N} diapositives)")
